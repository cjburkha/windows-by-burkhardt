"""
restyle_deck.py — Apex Energy Group brand restyle for Final_Drip_Campaign_Deck.pptx

Color palette from windowsbyburkhardt.com:
  Background:  #000000 (apex-black)
  Surface:     #1a1a1a (apex-dark-2)
  Accent:      #76bd1d (apex-green)
  Heading:     #ffffff
  Body:        #aaaaaa
  Border:      #2a2a2a

Run:
  python3 restyle_deck.py
Output:
  Final_Drip_Campaign_Deck_Apex.pptx
"""

import io
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree
import copy

# ── Brand colors ──────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK       = RGBColor(0x11, 0x11, 0x11)
SURFACE    = RGBColor(0x1a, 0x1a, 0x1a)
BORDER     = RGBColor(0x2a, 0x2a, 0x2a)
GREEN      = RGBColor(0x76, 0xbd, 0x1d)
WHITE      = RGBColor(0xff, 0xff, 0xff)
BODY       = RGBColor(0xaa, 0xaa, 0xaa)
MUTED      = RGBColor(0x66, 0x66, 0x66)

SRC  = '/Users/chrisb/dev/windows-by-burkhardt/gallary/Final_Drip_Campaign_Deck.pptx'
DEST = '/Users/chrisb/dev/windows-by-burkhardt/gallary/Final_Drip_Campaign_Deck_Apex.pptx'

# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width_pt=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width_pt:
            shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name='Calibri', font_size=18, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_window_logo(slide, left, top, size=0.55):
    """Draw the Apex window icon (window frame + cross dividers) as native shapes."""
    s = size
    # Outer border rect (stroke only)
    r = add_rect(slide, left, top, s, s, fill_color=None, line_color=GREEN, line_width_pt=2.5)
    # Horizontal bar
    bar_h = slide.shapes.add_shape(1,
        Inches(left), Inches(top + s/2 - 0.01),
        Inches(s), Inches(0.02))
    bar_h.fill.solid(); bar_h.fill.fore_color.rgb = GREEN
    bar_h.line.fill.background()
    # Vertical bar
    bar_v = slide.shapes.add_shape(1,
        Inches(left + s/2 - 0.01), Inches(top),
        Inches(0.02), Inches(s))
    bar_v.fill.solid(); bar_v.fill.fore_color.rgb = GREEN
    bar_v.line.fill.background()


def add_green_accent_bar(slide, left=0.4, top=1.15, width=0.05, height=0.45):
    """Thin vertical green accent bar next to section title."""
    add_rect(slide, left, top, width, height, fill_color=GREEN)


def add_footer(slide, label='Drip Campaign Proposal', slide_w=10.0):
    """Add a footer bar with label text."""
    add_rect(slide, 0, 7.2, slide_w, 0.3, fill_color=SURFACE)
    add_textbox(slide, label,
                0.3, 7.2, slide_w - 0.6, 0.3,
                font_size=8, color=MUTED, align=PP_ALIGN.LEFT)


def add_bullet_block(slide, bullets: list[str], left, top, width, height,
                     icon='•', icon_color=GREEN):
    """Add a styled bullet list box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        # Icon run (green)
        r1 = p.add_run()
        r1.text = icon + '  '
        r1.font.name = 'Calibri'
        r1.font.size = Pt(13)
        r1.font.color.rgb = icon_color
        r1.font.bold = True
        # Text run (body)
        r2 = p.add_run()
        r2.text = bullet
        r2.font.name = 'Calibri'
        r2.font.size = Pt(13)
        r2.font.color.rgb = BODY


def add_pros_cons_card(slide, left, top, width, height, title, items, title_color=GREEN, surface=SURFACE):
    """A card with a title and list of items."""
    # Card background
    card = add_rect(slide, left, top, width, height, fill_color=surface, line_color=BORDER, line_width_pt=1)
    # Title
    add_textbox(slide, title, left + 0.15, top + 0.12, width - 0.3, 0.35,
                font_size=13, bold=True, color=title_color)
    # Items
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.5),
        Inches(width - 0.3), Inches(height - 0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = '› ' + item
        r.font.name = 'Calibri'
        r.font.size = Pt(11)
        r.font.color.rgb = BODY


# ── Slide builders ────────────────────────────────────────────────────────────

def build_title_slide(slide):
    set_bg(slide, BLACK)

    # Large diagonal green accent stripe (decorative rect, rotated via XML)
    accent = add_rect(slide, 6.5, -1.0, 5.0, 10.0, fill_color=RGBColor(0x0d, 0x1f, 0x04))
    # Rotate 15 degrees
    sp = accent._element
    spPr = sp.find(qn('p:spPr'))
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    xfrm.set('rot', str(int(15 * 60000)))

    # Window logo — large, centered left half
    add_window_logo(slide, left=1.0, top=2.6, size=1.2)

    # Brand name
    add_textbox(slide, 'Apex Energy Group',
                2.5, 2.55, 5.5, 0.55,
                font_size=13, color=GREEN, bold=True)

    # Main title
    add_textbox(slide, 'Lead Generation\nReactivation Campaign',
                2.5, 3.1, 6.8, 1.4,
                font_size=32, bold=True, color=WHITE)

    # Subtitle
    add_textbox(slide, 'Milwaukee  •  Madison  •  Seattle',
                2.5, 4.55, 6.0, 0.4,
                font_size=13, color=BODY)

    add_textbox(slide, '4-Week Pilot Program',
                2.5, 4.95, 4.0, 0.35,
                font_size=11, italic=True, color=MUTED)

    # Green bottom stripe
    add_rect(slide, 0, 7.0, 10.0, 0.5, fill_color=GREEN)
    add_textbox(slide, 'Drip Campaign Proposal  |  Confidential',
                0.3, 7.05, 9.0, 0.4,
                font_size=9, color=BLACK, bold=True)


def build_objectives_slide(slide):
    set_bg(slide, BLACK)
    add_rect(slide, 0, 0, 10.0, 1.05, fill_color=DARK)
    add_window_logo(slide, left=0.25, top=0.22, size=0.55)
    add_textbox(slide, 'Campaign Objectives',
                1.05, 0.22, 7.0, 0.6,
                font_size=26, bold=True, color=WHITE)
    add_green_accent_bar(slide, left=0.95, top=0.22, width=0.05, height=0.55)

    objectives = [
        'Reactivate dormant leads into active opportunities',
        'Create consistent pipeline generation across all markets',
        'Segment campaigns: Existing, Canceled, and Prospect leads',
        'Drive qualified consultations through automated follow-up',
        'Build brand recall and trust over a structured 4-week cadence',
    ]
    add_bullet_block(slide, objectives, left=0.7, top=1.3, width=8.6, height=4.5, icon='▶')

    # Stats bar
    stats = [('3', 'Markets'), ('4', 'Weeks'), ('5', 'Lead Segments'), ('∞', 'Scalable')]
    for i, (num, label) in enumerate(stats):
        x = 0.5 + i * 2.35
        add_rect(slide, x, 5.9, 2.1, 1.0, fill_color=SURFACE, line_color=BORDER, line_width_pt=1)
        add_textbox(slide, num, x, 5.95, 2.1, 0.55,
                    font_size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, x, 6.55, 2.1, 0.3,
                    font_size=10, color=BODY, align=PP_ALIGN.CENTER)

    add_footer(slide)


def build_requirements_slide(slide):
    set_bg(slide, BLACK)
    add_rect(slide, 0, 0, 10.0, 1.05, fill_color=DARK)
    add_window_logo(slide, left=0.25, top=0.22, size=0.55)
    add_textbox(slide, 'Campaign Requirements',
                1.05, 0.22, 7.0, 0.6,
                font_size=26, bold=True, color=WHITE)
    add_green_accent_bar(slide, left=0.95, top=0.22, width=0.05, height=0.55)

    reqs = [
        ('🌐', 'Landing Pages', 'Localized pages for Milwaukee, Madison, and Seattle'),
        ('✉', 'Email',         'Sent from Managing Sales Partner — personalized per market'),
        ('💬', 'SMS',           'Localized reply routing to the correct regional team'),
        ('📊', 'Analytics',     'Full tracking: open rates, clicks, conversion per segment'),
        ('🔄', 'Automation',    'Triggered sequences based on lead status and engagement'),
    ]

    for i, (icon, title, desc) in enumerate(reqs):
        y = 1.3 + i * 1.1
        add_rect(slide, 0.5, y, 9.0, 0.95, fill_color=SURFACE, line_color=BORDER, line_width_pt=1)
        add_textbox(slide, icon, 0.65, y + 0.15, 0.5, 0.6, font_size=20, color=GREEN)
        add_textbox(slide, title, 1.3, y + 0.06, 2.2, 0.38, font_size=13, bold=True, color=WHITE)
        add_textbox(slide, desc,  1.3, y + 0.44, 7.8, 0.38, font_size=11, color=BODY)

    add_footer(slide)


def build_platform_slide(slide, title, strengths, limitations, cost):
    set_bg(slide, BLACK)
    add_rect(slide, 0, 0, 10.0, 1.05, fill_color=DARK)
    add_window_logo(slide, left=0.25, top=0.22, size=0.55)
    add_textbox(slide, 'Platform Option',
                1.05, 0.18, 4.0, 0.38,
                font_size=11, color=GREEN, bold=True)
    add_textbox(slide, title,
                1.05, 0.52, 7.5, 0.5,
                font_size=24, bold=True, color=WHITE)
    add_green_accent_bar(slide, left=0.95, top=0.18, width=0.05, height=0.85)

    # Strengths card
    add_pros_cons_card(slide, left=0.4, top=1.3, width=4.3, height=4.5,
                       title='✓  Strengths', items=strengths,
                       title_color=GREEN)

    # Limitations card
    add_pros_cons_card(slide, left=5.0, top=1.3, width=4.6, height=4.5,
                       title='✗  Limitations', items=limitations,
                       title_color=RGBColor(0xe0, 0x52, 0x52))

    # Cost badge
    add_rect(slide, 0.4, 6.0, 9.2, 0.7, fill_color=RGBColor(0x0d, 0x1f, 0x04), line_color=GREEN, line_width_pt=1.5)
    add_textbox(slide, '💰  Estimated Cost:  ' + cost,
                0.6, 6.05, 8.8, 0.55,
                font_size=14, bold=True, color=GREEN)

    add_footer(slide)


def build_recommendation_slide(slide):
    set_bg(slide, BLACK)
    add_rect(slide, 0, 0, 10.0, 1.05, fill_color=DARK)
    add_window_logo(slide, left=0.25, top=0.22, size=0.55)
    add_textbox(slide, 'Recommendation',
                1.05, 0.22, 7.0, 0.6,
                font_size=26, bold=True, color=WHITE)
    add_green_accent_bar(slide, left=0.95, top=0.22, width=0.05, height=0.55)

    steps = [
        ('01', 'Launch GoHighLevel Pilot',      'Start with the recommended all-in-one platform to validate the campaign model.'),
        ('02', 'Canceled Lead Campaign First',   'Target highest-ROI segment first — these leads already know the brand.'),
        ('03', 'Measure ROI + Conversion',       'Track results weekly: open rates, replies, booked consultations.'),
        ('04', 'Add Facebook Ads Integration',   'Layer paid retargeting once organic sequences are optimized.'),
        ('05', 'Expand to Additional Markets',   'Roll out the proven playbook to Madison, Seattle, and beyond.'),
    ]

    for i, (num, title, desc) in enumerate(steps):
        y = 1.3 + i * 1.05
        # Number badge
        add_rect(slide, 0.4, y + 0.08, 0.55, 0.55, fill_color=GREEN)
        add_textbox(slide, num, 0.4, y + 0.08, 0.55, 0.55,
                    font_size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, 1.1, y + 0.04, 4.5, 0.38,
                    font_size=13, bold=True, color=WHITE)
        add_textbox(slide, desc, 1.1, y + 0.42, 8.4, 0.38,
                    font_size=11, color=BODY)

    add_footer(slide)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    # Start fresh from the original (preserve slide count / layout metadata)
    prs = Presentation(SRC)

    # Clear all existing shapes from every slide
    for slide in prs.slides:
        sp_tree = slide.shapes._spTree
        for sp in list(sp_tree):
            tag = sp.tag.split('}')[-1] if '}' in sp.tag else sp.tag
            if tag not in ('sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp', 'grpSpPr', 'nvGrpSpPr'):
                continue
            sp_tree.remove(sp)

    slides = list(prs.slides)

    build_title_slide(slides[0])
    build_objectives_slide(slides[1])
    build_requirements_slide(slides[2])
    build_platform_slide(slides[3],
        title='GoHighLevel  —  Recommended',
        strengths=['All-in-one CRM + SMS + Email', 'Strong automation workflows', 'Scales across all markets', 'Native landing page builder'],
        limitations=['Less polished email design', 'Setup complexity', 'Learning curve for new users'],
        cost='$100 – $300 / month')
    build_platform_slide(slides[4],
        title='Mailchimp',
        strengths=['Best-in-class email design', 'Fast deployment', 'Strong analytics + reporting'],
        limitations=['Limited SMS capabilities', 'Requires integrations for CRM', 'Not a full CRM solution'],
        cost='$50 – $300 / month')
    build_platform_slide(slides[5],
        title='Custom Platform',
        strengths=['Fully tailored to Apex workflows', 'Complete data ownership', 'Maximum flexibility'],
        limitations=['High upfront development cost', 'Long build time (3–6 months)', 'Ongoing maintenance required'],
        cost='$10,000 – $50,000+')
    build_recommendation_slide(slides[6])

    prs.save(DEST)
    print(f'✅  Saved: {DEST}')


if __name__ == '__main__':
    main()
